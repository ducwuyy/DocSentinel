"""
Document parsing: multi-format to LLM-readable (Markdown/JSON).
PRD §5.2.5; docs/01 — PyMuPDF, python-docx, openpyxl, python-pptx.
"""

import io
from pathlib import Path

from app.models.parser import ParsedDocument, ParsedDocumentMetadata


def _safe_filename(name: str) -> str:
    """Sanitise filename for metadata (PRD §7.2 APP-01)."""
    return Path(name).name[:255] if name else "unknown"


def _parse_pdf(content: bytes, filename: str) -> ParsedDocument:
    import fitz  # PyMuPDF

    doc = fitz.open(stream=content, filetype="pdf")
    parts = []
    for page in doc:
        parts.append(page.get_text())
    doc.close()
    text = "\n\n".join(parts).strip() or "(no text extracted)"
    return ParsedDocument(
        format="markdown",
        content=text,
        metadata=ParsedDocumentMetadata(
            filename=_safe_filename(filename), type="application/pdf", pages=len(parts)
        ),
    )


def _parse_docx(content: bytes, filename: str) -> ParsedDocument:
    from docx import Document

    doc = Document(io.BytesIO(content))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text.strip() for cell in row.cells))
    text = "\n\n".join(parts).strip() or "(no text extracted)"
    return ParsedDocument(
        format="markdown",
        content=text,
        metadata=ParsedDocumentMetadata(
            filename=_safe_filename(filename),
            type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ),
    )


def _parse_xlsx(content: bytes, filename: str) -> ParsedDocument:
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"## {sheet.title}\n")
        for row in sheet.iter_rows(values_only=True):
            parts.append(" | ".join(str(c) if c is not None else "" for c in row))
    wb.close()
    text = "\n".join(parts).strip() or "(no content)"
    return ParsedDocument(
        format="markdown",
        content=text,
        metadata=ParsedDocumentMetadata(
            filename=_safe_filename(filename),
            type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ),
    )


def _parse_pptx(content: bytes, filename: str) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts = []
    for i, slide in enumerate(prs.slides):
        parts.append(f"## Slide {i + 1}\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
    text = "\n\n".join(parts).strip() or "(no text extracted)"
    return ParsedDocument(
        format="markdown",
        content=text,
        metadata=ParsedDocumentMetadata(
            filename=_safe_filename(filename),
            type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
    )


def _parse_plain(content: bytes, filename: str, content_type: str) -> ParsedDocument:
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        text = content.decode("latin-1", errors="replace")
    return ParsedDocument(
        format="markdown",
        content=text.strip() or "(empty)",
        metadata=ParsedDocumentMetadata(
            filename=_safe_filename(filename), type=content_type or "text/plain"
        ),
    )


# Allowed extensions and MIME (PRD §7.2 APP-01)
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md"}
EXTENSION_TO_PARSER = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".xlsx": _parse_xlsx,
    ".pptx": _parse_pptx,
    ".txt": lambda c, f: _parse_plain(c, f, "text/plain"),
    ".md": lambda c, f: _parse_plain(c, f, "text/markdown"),
}


def parse_file(content: bytes, filename: str) -> ParsedDocument:
    """
    Parse a single file into ParsedDocument.
    Raises ValueError for unsupported type or parse error.
    """
    path = Path(filename)
    suffix = path.suffix.lower()
    if suffix not in ALLOWED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type: {suffix}. Allowed: {sorted(ALLOWED_EXTENSIONS)}"
        )
    parser_fn = EXTENSION_TO_PARSER[suffix]
    return parser_fn(content, filename)
